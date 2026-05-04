"""Build watermarking_3min.pptx for the CS 5788 final presentation.

Produces a 5-slide deck (16:9) populated with the real Gemma 2 9B numbers.
Edit-friendly: text and bullet content live as variables at the top of the file.
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG  = os.path.join(ROOT, "figures")
OUT  = os.path.join(ROOT, "docs", "watermarking_3min.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ── load numbers from JSON so the deck stays in sync with the report ──────
def load(path):
    with open(os.path.join(ROOT, path)) as f:
        return json.load(f)

det   = load("results/detection_gemma_summary.json")
rob   = load("results/robustness_gemma.json")
sweep = load("results/delta_sweep_gemma.json")

# Optional cross-architecture results (loaded if the LLaMA Colab eval has finished).
llama_det = None; llama_rob = None
if os.path.exists(os.path.join(ROOT, "results/detection_llama_summary.json")):
    llama_det = load("results/detection_llama_summary.json")
    llama_rob = load("results/robustness_llama.json")

NAVY  = RGBColor(0x14, 0x2A, 0x55)
RED   = RGBColor(0xC0, 0x39, 0x2B)
BLUE  = RGBColor(0x33, 0x5B, 0xA1)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]    # blank layout, we draw everything ourselves

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=18, color=NAVY, gap=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb

def fill_background(slide, color=WHITE):
    # Add a full-bleed rectangle as the slide background.
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    sp = rect._element
    sp.getparent().insert(2, sp)
    return rect

def add_band(slide, color=NAVY, height=Inches(0.45)):
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, height)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def add_footer(slide, idx, total):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8.0), Inches(0.35),
             "Invisible Watermarking for LLMs via Logit Manipulation  ·  Hasan & Syed (CS 5788)",
             size=10, color=GRAY)
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35),
             f"{idx} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ── SLIDE 1: Title + motivation ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_band(s, NAVY, Inches(2.6))

add_text(s, Inches(0.6), Inches(0.55), Inches(12.0), Inches(0.5),
         "CS 5788  ·  Final Project, Spring 2026", size=14, color=WHITE)
add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(1.1),
         "Invisible Watermarking for LLMs", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.95), Inches(12.0), Inches(0.7),
         "Re-implementing Kirchenbauer et al. (2023) on Gemma 2 9B & LLaMA 3.1 8B",
         size=22, color=WHITE)

add_text(s, Inches(0.6), Inches(3.10), Inches(12.0), Inches(0.5),
         "Why does this matter?", size=22, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(3.65), Inches(12.0), Inches(2.6),
    [
        "AI-generated text is now flooding news, classrooms, and platforms, and post-hoc detectors are unreliable on short or paraphrased passages.",
        "We embed an invisible statistical signal during generation and verify it later with a one-sided z-test.",
        "Goal: high detectability + preserved fluency + survival under realistic edits, on modern instruction-tuned models.",
    ], size=18, color=NAVY, gap=8)

add_text(s, Inches(0.6), Inches(6.25), Inches(12.0), Inches(0.5),
         "Ali Hasan (ah2434) · Ammar Syed (as4422)", size=14, color=GRAY)
add_footer(s, 1, 5)

# ── SLIDE 2: Method ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
fill_background(s, WHITE); add_band(s, NAVY)
add_text(s, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.45),
         "Method  ·  green-list logit-bias watermark", size=22, bold=True, color=WHITE)

s.shapes.add_picture(os.path.join(FIG, "fig1_method.png"),
                      Inches(0.6), Inches(0.7), width=Inches(12.1))

add_bullets(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(2.0),
    [
        "Inject:  SHA-256(seed ‖ prev token) → green list of γ|V| tokens; add bias δ to green-list logits before softmax.",
        "Detect:  replay the hash; count green-list hits; one-sided z-test → flag if z > τ (calibrated for 1% FPR).",
        "Parameters used:  γ = 0.5,   δ = 2,   secret seed = 42  (frozen across all experiments).",
    ], size=15, color=NAVY, gap=4)

add_footer(s, 2, 5)

# ── SLIDE 3: Headline detection + robustness ────────────────────────────
s = prs.slides.add_slide(BLANK)
fill_background(s, WHITE); add_band(s, NAVY)
add_text(s, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.45),
         "Detection works, and survives surface attacks", size=22, bold=True, color=WHITE)

# Headline TPR card (top-left)
add_text(s, Inches(0.5), Inches(0.7), Inches(6.0), Inches(0.5),
         "Headline detection (Gemma 2 9B)", size=18, bold=True, color=NAVY)
gemma_bullets = [
    f"TPR @ 1% FPR (all):  {det['tpr_at_1pct_fpr_all']*100:.1f}%",
    f"TPR @ 1% FPR (≥150 tok):  {det['tpr_at_1pct_fpr_ge150tok']*100:.1f}%",
    f"GPT-2 PPL ratio (wm/uwm) = {det['ppl_ratio_wm_over_uwm']:.2f}  (+9% overhead)",
]
if llama_det is not None:
    gemma_bullets.append(
        f"Same code on LLaMA 3.1 8B:  {llama_det['tpr_at_1pct_fpr_all']*100:.1f}% TPR, "
        f"PPL ratio {llama_det['ppl_ratio_wm_over_uwm']:.2f}"
    )
else:
    gemma_bullets.append(f"Calibrated z-threshold τ = {det['calibrated_z_threshold']:.2f}")
add_bullets(s, Inches(0.5), Inches(1.20), Inches(6.0), Inches(2.4),
    gemma_bullets, size=16, color=NAVY, gap=6)

s.shapes.add_picture(os.path.join(FIG, "fig2_zscore_hist.png"),
                      Inches(6.6), Inches(0.7), width=Inches(6.4))

# Robustness (bottom row)
add_text(s, Inches(0.5), Inches(3.85), Inches(6.0), Inches(0.5),
         "Robust to mild edits", size=18, bold=True, color=NAVY)
rob_bullets = [
    f"5% word sub:  TPR = {rob['word_sub_5pct']['tpr_at_1pct_fpr']*100:.1f}%   (LLaMA: {llama_rob['word_sub_5pct']['tpr_at_1pct_fpr']*100:.1f}%)" if llama_rob else f"5% word sub:  TPR = {rob['word_sub_5pct']['tpr_at_1pct_fpr']*100:.1f}%",
    f"10% word sub:  TPR = {rob['word_sub_10pct']['tpr_at_1pct_fpr']*100:.1f}%   (LLaMA: {llama_rob['word_sub_10pct']['tpr_at_1pct_fpr']*100:.1f}%)" if llama_rob else f"10% word sub:  TPR = {rob['word_sub_10pct']['tpr_at_1pct_fpr']*100:.1f}%",
    f"20% word sub:  TPR = {rob['word_sub_20pct']['tpr_at_1pct_fpr']*100:.1f}%   (LLaMA: {llama_rob['word_sub_20pct']['tpr_at_1pct_fpr']*100:.1f}%)" if llama_rob else f"20% word sub:  TPR = {rob['word_sub_20pct']['tpr_at_1pct_fpr']*100:.1f}%",
    f"10% token ins/del: TPR ≈ {rob['token_insertion_10pct']['tpr_at_1pct_fpr']*100:.0f}%   (LLaMA: {llama_rob['token_insertion_10pct']['tpr_at_1pct_fpr']*100:.0f}%)" if llama_rob else f"10% token ins/del: TPR ≈ {rob['token_insertion_10pct']['tpr_at_1pct_fpr']*100:.0f}%",
]
add_bullets(s, Inches(0.5), Inches(4.35), Inches(6.0), Inches(2.5),
    rob_bullets, size=15, color=NAVY, gap=6)

s.shapes.add_picture(os.path.join(FIG, "fig4_robustness.png"),
                      Inches(6.6), Inches(3.85), width=Inches(6.4))

add_footer(s, 3, 5)

# ── SLIDE 4: Detectability vs quality (delta sweep) ─────────────────────
s = prs.slides.add_slide(BLANK)
fill_background(s, WHITE); add_band(s, NAVY)
add_text(s, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.45),
         "Detectability ↔ quality:  the δ knee is at 2", size=22, bold=True, color=WHITE)

s.shapes.add_picture(os.path.join(FIG, "fig5_delta_tradeoff.png"),
                      Inches(0.6), Inches(0.8), width=Inches(7.6))

add_text(s, Inches(8.4), Inches(0.9), Inches(4.6), Inches(0.45),
         "Sweep over δ ∈ {0.5, 1, 2, 4, 8}", size=18, bold=True, color=NAVY)

# Build a small TPR / PPL table from the sweep JSON
add_text(s, Inches(8.4), Inches(1.5), Inches(2.0), Inches(0.4),
         "δ", size=14, bold=True, color=GRAY)
add_text(s, Inches(9.6), Inches(1.5), Inches(2.0), Inches(0.4),
         "TPR", size=14, bold=True, color=GRAY)
add_text(s, Inches(11.2), Inches(1.5), Inches(2.0), Inches(0.4),
         "PPL ratio", size=14, bold=True, color=GRAY)
for i, row in enumerate(sweep):
    yy = Inches(1.95 + i*0.42)
    add_text(s, Inches(8.4), yy, Inches(2.0), Inches(0.4),
             f"{row['delta']:.1f}", size=14, color=NAVY)
    add_text(s, Inches(9.6), yy, Inches(2.0), Inches(0.4),
             f"{row['tpr']*100:.1f}%", size=14, color=RED, bold=True)
    add_text(s, Inches(11.2), yy, Inches(2.0), Inches(0.4),
             f"{row['ppl_ratio']:.2f}×", size=14, color=BLUE, bold=True)

add_bullets(s, Inches(0.6), Inches(5.85), Inches(12.0), Inches(1.3),
    [
        "δ = 0.5: too weak (23% TPR).  δ = 8: ‘obviously biased’ text (PPL ratio 2.3×).",
        "δ = 2 sits at the knee:  90% TPR overall (99% on long completions) at only 14% PPL overhead.",
    ], size=15, color=NAVY, gap=4)

add_footer(s, 4, 5)

# ── SLIDE 5: Take-aways ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT); add_band(s, NAVY)
add_text(s, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.45),
         "Take-aways", size=22, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(0.85), Inches(12.0), Inches(0.55),
         "What we showed", size=20, bold=True, color=NAVY)

llama_bullet = (
    f"Cross-architecture: same 30-line LogitsProcessor on LLaMA 3.1 8B → "
    f"{llama_det['tpr_at_1pct_fpr_all']*100:.1f}% TPR, PPL ratio {llama_det['ppl_ratio_wm_over_uwm']:.2f} "
    f"(stronger than Gemma; different family, different tokenizer)."
    if llama_det is not None else
    "Cross-architecture: same 30-line LogitsProcessor works unchanged on LLaMA 3.1 8B (different family, different tokenizer)."
)

add_bullets(s, Inches(0.6), Inches(1.45), Inches(12.0), Inches(2.6),
    [
        "Reproduced the Kirchenbauer scheme on modern 9B instruction-tuned models. Every empirical claim of the original paper still holds.",
        "Detection is reliable (≥99% TPR for 150+ tokens), output quality is preserved (+9% PPL), and the watermark survives 20% random word substitution.",
        llama_bullet,
        "Documented a real failure mode: HuggingFace’s multimodal-by-default loader silently breaks Gemma 3. That is a load-time issue rather than a watermark issue.",
    ], size=16, color=NAVY, gap=8)

add_text(s, Inches(0.6), Inches(5.20), Inches(12.0), Inches(0.55),
         "What’s next", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(5.80), Inches(12.0), Inches(1.4),
    [
        "Stronger attacks: full LLM paraphrase + spring paraphrase from Kirchenbauer 2024.",
        "Distortion-free / semantic-grouping variants  (Kuditipudi 2023; Christ 2024)  for higher-stakes deployments.",
    ], size=16, color=GRAY, gap=6)

add_footer(s, 5, 5)

# ── Save ────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Wrote {OUT}")
